import glob
import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from data_analysis.utils import Utils

utils = Utils(
    is_normal=True,
    is_previous=False,
    data_type="spectrogram",
    fit_pos="middle",
    stride="middle",  # FIXME: 要チェック
    kernel_size=128,
    model_type="enn",
    cleansing_type="noting",
)
tmp_dir = utils.env.tmp_dir
filepath_list = [
    os.path.join(root_dir, filepath[0])
    for root_dir, _, filepath in os.walk(tmp_dir)
    if len(filepath) != 0
]
metrics_filepath_list = [
    os.path.join(root_dir, filepath[1])
    for root_dir, _, filepath in os.walk(tmp_dir)
    if len(filepath) != 0
]
cm_filepath_list = [
    os.path.join(root_dir, filepath[2])
    for root_dir, _, filepath in os.walk(tmp_dir)
    if len(filepath) != 0
]

pptx_filepath = (
    "/home/takadamalab/taiki_senju/git/sleep_study/gallery/blue_template.pptx"
)
prs = Presentation(pptx_filepath)

title_slide_layout = prs.slide_layouts[3]
for ss_filepath, cm_filepath, metrics in zip(
    filepath_list, cm_filepath_list, metrics_filepath_list
):
    slide = prs.slides.add_slide(title_slide_layout)
    left = Inches(0)
    top = Inches(0)
    pic = slide.shapes.add_picture(ss_filepath, left, top)
    left = Inches(7)
    pic = slide.shapes.add_picture(cm_filepath, left, top)
    left, top, width, height = Inches(8), Inches(5), Inches(4), Inches(2)
    df = pd.read_csv(metrics)
    __row = int((df.shape[1] - 1) / 3)
    table = slide.shapes.add_table(3, __row, left, top, width, height).table
    df_ndarray = df.to_numpy()[:, 1:]
    for i, row in enumerate(df_ndarray):
        for j in range(__row):
            table.cell(i, j).text = f"{row[j + __row * i]:.2f}"


prs.save(os.path.join(os.path.split(pptx_filepath)[0], "output.pptx"))
